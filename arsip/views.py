from django.shortcuts import render, redirect, get_object_or_404
from django.urls import reverse_lazy
from django.views.generic import CreateView, ListView, UpdateView
from django.utils import timezone
from .models import Arsip, Kategori, Profile
from .forms import ArsipForm, VerifikasiForm, KategoriForm, RegisterForm, PasswordChangeForm
from django.core.paginator import Paginator
from django.db.models import Count, Q
from rest_framework import viewsets
from django_filters.rest_framework import DjangoFilterBackend
from .serializers import ArsipSerializer
from django.contrib.auth.decorators import login_required, user_passes_test
from django.contrib.auth.mixins import LoginRequiredMixin
from .permissions import PembuatRequiredMixin, VerifikasiRequiredMixin
from django.contrib.auth import login
from django.contrib.auth.models import User
from django.contrib import messages
import mimetypes
from django.http import HttpResponse, JsonResponse
from django.views.decorators.clickjacking import xframe_options_exempt
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import io
import os

# Helper function to check if a user is Admin or Editor
def is_admin_or_editor(user):
    return user.is_authenticated and hasattr(user, 'profile') and user.profile.role in ['Admin', 'Editor']

# Helper function to check if a user is Admin or the owner of the archive
def is_admin_or_arsip_owner(user, arsip):
    return user.is_authenticated and hasattr(user, 'profile') and (user.profile.role == 'Admin' or user == arsip.dibuat_oleh)

# Helper function to increment archive access
def increment_arsip_access(request, arsip_id):
    """
    Increments the 'diakses' column in Arsip and tracks user session
    to prevent double counting within one session.
    """
    if 'accessed_archives' not in request.session:
        request.session['accessed_archives'] = []

    if arsip_id not in request.session['accessed_archives']:
        try:
            arsip = Arsip.objects.get(id=arsip_id)
            arsip.diakses = arsip.diakses + 1
            arsip.save(update_fields=['diakses'])

            request.session['accessed_archives'].append(arsip_id)
            request.session.modified = True
            print(f"Access for Archive ID {arsip_id} incremented. Current access count: {arsip.diakses}")
        except Arsip.DoesNotExist:
            print(f"Archive ID {arsip_id} not found.")
        except Exception as e:
            print(f"Failed to increment access for Archive ID {arsip_id}: {e}")
    else:
        print(f"Archive ID {arsip_id} already accessed in this session. Not incrementing access.")


# View function for document preview
@xframe_options_exempt
def view_digital_document(request, arsip_id):
    increment_arsip_access(request, arsip_id)

    try:
        arsip = get_object_or_404(Arsip, id=arsip_id)
        file_path = arsip.lokasi_digital.path
        file_name = os.path.basename(file_path)
        file_extension = os.path.splitext(file_path)[1].lower()

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found at {file_path}")

        if arsip.status_akses == 'internal' and not request.user.is_authenticated:
             return JsonResponse({'error': 'Akses ditolak. File ini hanya untuk pengguna internal.', 'download_url': None}, status=403)


        if file_extension == '.pdf':
            with open(file_path, 'rb') as f:
                response = HttpResponse(f.read(), content_type='application/pdf')
                response['Content-Disposition'] = f'inline; filename="{file_name}"'
                return response
        elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
            with open(file_path, 'rb') as f:
                content_type, _ = mimetypes.guess_type(file_name)
                response = HttpResponse(f.read(), content_type=content_type or 'application/octet-stream')
                response['Content-Disposition'] = f'inline; filename="{file_name}"'
                return response
        elif file_extension == '.docx':
            doc = Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return HttpResponse("\n".join(full_text), content_type='text/plain')
        elif file_extension == '.xlsx':
            wb = load_workbook(file_path)
            sheet = wb.active
            text_content = []
            for row in sheet.iter_rows():
                row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                text_content.append("\t".join(row_values))
            return HttpResponse("\n".join(text_content), content_type='text/plain')
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            return HttpResponse("\n".join(text_content), content_type='text/plain')
        else:
            return JsonResponse({'error': 'Jenis file tidak dapat ditampilkan langsung, silahkan unduh.', 'download_url': arsip.lokasi_digital.url}, status=400)
    except FileNotFoundError:
        return JsonResponse({'error': 'File tidak ditemukan.'}, status=404)
    except Exception as e:
        print(f"Error processing file for view: {e}")
        return JsonResponse({'error': f'Terjadi kesalahan saat memproses file: {str(e)}'}, status=500)


# View function for document download
@login_required
def download_digital_document(request, arsip_id):
    increment_arsip_access(request, arsip_id)

    try:
        arsip = get_object_or_404(Arsip, id=arsip_id)
        file_path = arsip.lokasi_digital.path
        file_name = os.path.basename(file_path)

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found at {file_path}")

        if arsip.status_akses == 'internal' and not request.user.is_authenticated:
            messages.error(request, 'Anda tidak memiliki akses untuk mengunduh arsip internal ini. Silakan login.')
            return redirect('arsip:login')

        with open(file_path, 'rb') as f:
            content_type = mimetypes.guess_type(file_name)[0] or 'application/octet-stream'
            response = HttpResponse(f.read(), content_type=content_type)
            response['Content-Disposition'] = f'attachment; filename="{file_name}"'
            return response
    except FileNotFoundError:
        messages.error(request, "File yang diminta tidak ditemukan.")
        return redirect('arsip:detail_arsip', arsip_id=arsip_id)
    except Exception as e:
        messages.error(request, f"Terjadi kesalahan saat mengunduh file: {str(e)}")
        return redirect('arsip:detail_arsip', arsip_id=arsip_id)


@login_required
def dashboard(request):
    profile = request.user.profile

    if profile.role == 'Admin':
        # Perlu didefinisikan di urls.py Anda
        return redirect('dashboard_admin')
    elif profile.role == 'Editor':
        return redirect('dashboard_editor')
    elif profile.role == 'Pembuat':
        return redirect('dashboard_pembuat')
    else:
        return redirect('arsip:beranda')


@login_required
def dashboard_pembuat(request):
    # Memastikan pengguna memiliki atribut 'profile' sebelum mengaksesnya
    if not hasattr(request.user, 'profile') or request.user.profile.role != 'Pembuat':
        messages.error(request, 'Anda tidak memiliki izin untuk mengakses halaman ini.')
        return redirect('arsip:beranda')

    arsip_pribadi = Arsip.objects.filter(
        dibuat_oleh=request.user
    ).order_by('-tanggal_diterima')

    return render(request, 'arsip/dashboard_pembuat.html', {
        'arsip_pribadi': arsip_pribadi
    })

@login_required
def dashboard_editor(request):
    # Memastikan pengguna memiliki atribut 'profile' sebelum mengaksesnya
    if not hasattr(request.user, 'profile') or request.user.profile.role not in ['Admin', 'Editor']:
        messages.error(request, 'Anda tidak memiliki izin untuk mengakses halaman ini.')
        return redirect('arsip:beranda')

    arsip_pending = Arsip.objects.filter(status_verifikasi='pending').count()
    arsip_diverifikasi = Arsip.objects.filter(
        Q(verifikator=request.user) | Q(status_verifikasi='diterima')
    ).order_by('-tanggal_verifikasi')[:5]

    return render(request, 'arsip/dashboard_editor.html', {
        'arsip_pending': arsip_pending,
        'arsip_diverifikasi': arsip_diverifikasi
    })

class ArsipViewSet(viewsets.ModelViewSet):
    queryset = Arsip.objects.all()
    serializer_class = ArsipSerializer
    filter_backends = [DjangoFilterBackend]
    filterset_fields = ['kategori', 'status_akses']

def beranda(request):
    # Mulai dengan arsip yang sudah diverifikasi
    arsip_terbaru_queryset = Arsip.objects.filter(status_verifikasi='diterima')

    # Jika pengguna tidak login, hanya tampilkan arsip publik
    if not request.user.is_authenticated:
        arsip_terbaru_queryset = arsip_terbaru_queryset.filter(status_akses='publik')

    arsip_terbaru = arsip_terbaru_queryset.order_by('-tanggal_diterima')[:6]

    kategori_populer = Kategori.objects.annotate(
        jumlah_arsip=Count('arsip')
    ).order_by('-jumlah_arsip')[:4]

    return render(request, 'arsip/beranda.html', {
        'arsip_terbaru': arsip_terbaru,
        'kategori_populer': kategori_populer,
    })

def daftar_arsip(request):
    # Mulai dengan semua arsip yang sudah diverifikasi
    arsip_list = Arsip.objects.filter(status_verifikasi='diterima')
    
    kategori_id = request.GET.get('kategori')
    status_akses_filter = request.GET.get('status_akses') 

    # Filter berdasarkan kategori jika ada
    if kategori_id:
        arsip_list = arsip_list.filter(kategori_id=kategori_id)
    
    # Perbaikan: Awalnya filter ini di dalam `if request.user.is_authenticated:`
    # Sekarang kita proses status_akses_filter terlebih dahulu.
    if status_akses_filter:
        if status_akses_filter == 'publik':
            arsip_list = arsip_list.filter(status_akses='publik')
        elif status_akses_filter == 'internal':
            # Pastikan user terautentikasi dan memiliki profil sebelum mengakses role
            if request.user.is_authenticated and hasattr(request.user, 'profile'):
                arsip_list = arsip_list.filter(status_akses='internal')
            else:
                # Jika user mencoba mengakses internal tapi tidak login/tidak punya profil
                messages.warning(request, "Anda harus login untuk melihat arsip internal.")
                arsip_list = arsip_list.none() # Atau bisa juga redirect ke login page
    else:
        # Default behavior ketika tidak ada filter status_akses:
        if not request.user.is_authenticated:
            # Jika tidak login, defaultnya hanya tampilkan publik
            arsip_list = arsip_list.filter(status_akses='publik')
        # Jika login dan tidak ada filter status_akses, biarkan arsip_list apa adanya (publik dan internal)

    arsip_list = arsip_list.order_by('-tanggal_diterima')
    paginator = Paginator(arsip_list, 10)
    page_number = request.GET.get('page')
    page_obj = paginator.get_page(page_number)

    kategori_list = Kategori.objects.all().order_by('nama')

    return render(request, 'arsip/daftar_arsip.html', {
        'page_obj': page_obj,
        'kategori_list': kategori_list,
    })

def detail_arsip(request, arsip_id):
    arsip = get_object_or_404(Arsip, id=arsip_id)

    if arsip.status_akses == 'internal' and not request.user.is_authenticated:
        messages.error(request, 'Anda tidak memiliki akses ke arsip internal ini. Silakan login.')
        return redirect('arsip:login')

    arsip_terkait = Arsip.objects.filter(kategori=arsip.kategori).exclude(id=arsip_id)[:3]
    return render(request, 'arsip/detail_arsip.html', {'arsip': arsip, 'arsip_terkait': arsip_terkait})

@login_required
@user_passes_test(is_admin_or_editor, login_url=reverse_lazy('arsip:beranda'))
def edit_arsip(request, arsip_id):
    arsip = get_object_or_404(Arsip, id=arsip_id)
    
    # Memastikan pengguna memiliki atribut 'profile' sebelum mengaksesnya
    if not (request.user.is_authenticated and hasattr(request.user, 'profile') and (request.user.profile.role == 'Admin' or request.user.profile.role == 'Editor' or request.user == arsip.dibuat_oleh)):
        messages.error(request, "Anda tidak memiliki izin untuk mengedit arsip ini!")
        return redirect('arsip:detail_arsip', arsip_id=arsip.id)

    if request.method == 'POST':
        form = ArsipForm(request.POST, request.FILES, instance=arsip)
        if form.is_valid():
            # Logika baru untuk penanganan file:
            # 1. Cek apakah ada file baru diupload
            if 'lokasi_digital' in request.FILES:
                # Jika ada file lama, hapus file lama secara fisik
                if arsip.lokasi_digital:
                    old_file_path = arsip.lokasi_digital.path
                    if os.path.exists(old_file_path):
                        os.remove(old_file_path)
                        print(f"File lama {old_file_path} berhasil dihapus.")
                # Atur file baru dari form
                arsip.lokasi_digital = form.cleaned_data['lokasi_digital']
            # 2. Cek apakah kotak "Clear" dicentang
            elif f'lokasi_digital-clear' in request.POST and form.cleaned_data.get('lokasi_digital') is False:
                if arsip.lokasi_digital:
                    old_file_path = arsip.lokasi_digital.path
                    if os.path.exists(old_file_path):
                        os.remove(old_file_path)
                        print(f"File {old_file_path} berhasil dihapus karena opsi 'clear'.")
                arsip.lokasi_digital = None # Set field file menjadi None (kosong)
            # 3. Jika tidak ada file baru diupload DAN tidak ada clear, biarkan file lama
            # (form.cleaned_data['lokasi_digital'] akan tetap berisi instance file lama dari `instance=arsip`)
            
            arsip = form.save(commit=False) # Simpan data form ke instance Arsip, tapi belum ke DB
            
            # Perbarui format_file dan ukuran_file jika lokasi_digital diubah atau dikosongkan
            if arsip.lokasi_digital:
                arsip.format_file = os.path.splitext(arsip.lokasi_digital.name)[1][1:].upper()
                arsip.ukuran_file = arsip.lokasi_digital.size
            else:
                arsip.format_file = ''
                arsip.ukuran_file = 0

            arsip.save() # Simpan instance Arsip ke DB
            form.save_m2m() # Simpan many-to-many data (misal tags)
            messages.success(request, 'Arsip berhasil diperbarui!')
            return redirect('arsip:detail_arsip', arsip_id=arsip.id)
        else:
            messages.error(request, "Ada kesalahan dalam formulir. Mohon periksa kembali.")
    else:
        form = ArsipForm(instance=arsip)
    
    return render(request, 'arsip/edit_arsip.html', {'form': form, 'arsip': arsip})


@login_required
def hapus_arsip(request, arsip_id):
    arsip = get_object_or_404(Arsip, id=arsip_id)
    
    # Memastikan pengguna memiliki atribut 'profile' sebelum mengaksesnya
    if not (request.user.is_authenticated and hasattr(request.user, 'profile') and (request.user.profile.role == 'Admin' or request.user == arsip.dibuat_oleh)):
        messages.error(request, "Anda tidak memiliki izin untuk menghapus arsip ini!")
        return redirect('arsip:detail_arsip', arsip_id=arsip.id)
    
    if request.method == 'POST':
        if arsip.lokasi_digital:
            file_path = arsip.lokasi_digital.path
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                    print(f"File {file_path} berhasil dihapus.")
                except Exception as e:
                    print(f"Gagal menghapus file {file_path}: {e}")
                    messages.error(request, f"Gagal menghapus file fisik: {e}")
                    return redirect('arsip:detail_arsip', arsip_id=arsip_id)
        
        arsip.delete()
        messages.success(request, "Arsip berhasil dihapus!")
        return redirect('arsip:daftar_arsip')
    
    return render(request, 'arsip/konfirmasi_hapus.html', {'arsip': arsip})

# Fungsi baru untuk menambah kategori
@login_required
@user_passes_test(lambda user: hasattr(user, 'profile') and user.profile.role == 'Admin', login_url=reverse_lazy('arsip:beranda'))
def tambah_kategori(request):
    if request.method == 'POST':
        form = KategoriForm(request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, 'Kategori baru berhasil ditambahkan!')
            return redirect('arsip:kategori_list') # Redirect ke daftar kategori setelah sukses
        else:
            messages.error(request, 'Terjadi kesalahan saat menambahkan kategori. Mohon periksa isian form.')
    else:
        form = KategoriForm()
    
    return render(request, 'arsip/tambah_kategori.html', {'form': form})

# New function to edit category
@login_required
# PERBAIKAN: Izinkan Admin ATAU Editor untuk mengedit kategori
@user_passes_test(lambda user: hasattr(user, 'profile') and user.profile.role in ['Admin', 'Editor'], login_url=reverse_lazy('arsip:beranda'))
def edit_kategori(request, kategori_id):
    kategori = get_object_or_404(Kategori, id=kategori_id)
    if request.method == 'POST':
        form = KategoriForm(request.POST, instance=kategori)
        if form.is_valid():
            form.save()
            messages.success(request, f'Kategori "{kategori.nama}" berhasil diperbarui.')
            return redirect('arsip:kategori_list')
        else:
            messages.error(request, 'Terjadi kesalahan saat memperbarui kategori. Mohon periksa isian form.')
    else:
        form = KategoriForm(instance=kategori)
    return render(request, 'arsip/edit_kategori.html', {'form': form, 'kategori': kategori})

# New function to delete category
@login_required
@user_passes_test(lambda user: hasattr(user, 'profile') and user.profile.role == 'Admin', login_url=reverse_lazy('arsip:beranda'))
def hapus_kategori(request, kategori_id):
    kategori = get_object_or_404(Kategori, id=kategori_id)

    if request.method == 'POST':
        # When a category is deleted, the 'kategori' field in related Archives will automatically become NULL
        # due to models.ForeignKey(Kategori, on_delete=models.SET_NULL)
        kategori.delete()
        messages.success(request, f'Kategori "{kategori.nama}" berhasil dihapus.')
        return redirect('arsip:kategori_list')
    
    # If method GET, show confirmation page (optional, could redirect directly if no confirmation is needed)
    # For security, a confirmation page or modal is better.
    return render(request, 'arsip/konfirmasi_hapus_kategori.html', {'kategori': kategori})

def kategori_populer(request):
    kategori_populer = Kategori.objects.annotate(
        jumlah_arsip=Count('arsip')
    ).order_by('-jumlah_arsip')[:4]

    return render(request, 'arsip/kategori_populer.html', {
        'kategori_populer': kategori_populer
    })

def kategori_list(request):
    kategori = Kategori.objects.all().annotate(
        jumlah_arsip=Count('arsip')
    ).order_by('nama')
    return render(request, 'arsip/kategori_list.html', {
        'kategori_list': kategori
    })

def arsip_per_kategori(request, kategori_id):
    kategori = get_object_or_404(Kategori, id=kategori_id)
    
    # Initial filter: Only verified archives
    arsip_list = Arsip.objects.filter(kategori=kategori, status_verifikasi='diterima')
    
    # Get access status filter from URL (e.g., ?status_akses=internal)
    status_akses_filter = request.GET.get('status_akses') 

    if status_akses_filter:
        if status_akses_filter == 'publik':
            arsip_list = arsip_list.filter(status_akses='publik')
        elif status_akses_filter == 'internal':
            # Only show internal if user is authenticated and has a profile
            if request.user.is_authenticated and hasattr(request.user, 'profile'):
                arsip_list = arsip_list.filter(status_akses='internal')
            else:
                messages.warning(request, "Anda harus login untuk melihat arsip internal di kategori ini.")
                arsip_list = arsip_list.none() # Empty queryset if no access
    else:
        # Changes here: If no status_akses filter in URL
        if not request.user.is_authenticated:
            # If not logged in, default to only public
            arsip_list = arsip_list.filter(status_akses='publik')
        # If logged in, and no status_akses filter, leave arsip_list as is (public and internal)

    arsip_list = arsip_list.order_by('-tanggal_diterima')

    return render(request, 'arsip/arsip_per_kategori.html', {
        'kategori': kategori,
        'arsip_list': arsip_list,
    })

def cari_arsip(request):
    query = request.GET.get('q', '')
    status_filter = request.GET.get('status_akses', 'publik') # Change from 'status' to 'status_akses'
    
    results = Arsip.objects.filter(
        Q(judul_arsip__icontains=query) |
        Q(deskripsi__icontains=query) |
        Q(kategori__nama__icontains=query)
    ).filter(status_verifikasi='diterima') # Ensure only verified

    if not request.user.is_authenticated:
        results = results.filter(status_akses='publik')
    else:
        if status_filter == 'internal':
            results = results.filter(status_akses='internal')
        elif status_filter == 'publik':
            results = results.filter(status_akses='publik')
        # If user is logged in and status_filter is empty (default), show all (public and internal)
            
    results = results.order_by('-tanggal_diterima')

    return render(request, 'arsip/hasil_pencarian.html', {
        'query': query,
        'results': results,
    })

class UploadArsipView(LoginRequiredMixin, PembuatRequiredMixin, CreateView):
    model = Arsip
    form_class = ArsipForm
    template_name = 'arsip/upload.html'
    success_url = reverse_lazy('arsip:daftar_arsip_pribadi')

    def form_valid(self, form):
        form.instance.dibuat_oleh = self.request.user
        form.instance.tanggal_diterima = timezone.now().date()
        form.instance.status_verifikasi = 'pending'
        return super().form_valid(form)

class VerifikasiArsipView(LoginRequiredMixin, VerifikasiRequiredMixin, ListView):
    model = Arsip
    template_name = 'arsip/verifikasi_list.html'
    context_object_name = 'arsip_list'
    paginate_by = 10

    def get_queryset(self):
        queryset = Arsip.objects.filter(status_verifikasi='pending').order_by('-tanggal_diterima')
        return queryset

    def dispatch(self, request, *args, **kwargs):
        if not request.user.is_authenticated or not hasattr(request.user, 'profile') or request.user.profile.role not in ['Admin', 'Editor']:
            messages.error(request, 'Anda tidak memiliki izin untuk mengakses halaman verifikasi.')
            return redirect('arsip:beranda')
        return super().dispatch(request, *args, **kwargs)

class VerifikasiDetailView(LoginRequiredMixin, VerifikasiRequiredMixin, UpdateView):
    model = Arsip
    form_class = VerifikasiForm
    template_name = 'arsip/verifikasi_detail.html'
    success_url = reverse_lazy('arsip:verifikasi_arsip')
    context_object_name = 'object'

    def form_valid(self, form):
        form.instance.verifikator = self.request.user
        form.instance.tanggal_verifikasi = timezone.now()
        response = super().form_valid(form)
        self.object.save()
        return response

    def dispatch(self, request, *args, **kwargs):
        if not request.user.is_authenticated or not hasattr(request.user, 'profile') or request.user.profile.role not in ['Admin', 'Editor']:
            messages.error(request, 'Anda tidak memiliki izin untuk mengakses halaman verifikasi.')
            return redirect('arsip:beranda')
        return super().dispatch(request, *args, **kwargs)

@login_required
def profil(request):
    return render(request, 'arsip/profil.html', {
        'user': request.user,
        'profile': request.user.profile,
    })

def statistik_context(request):
    from .models import Arsip, Kategori
    return {
        'total_arsip': Arsip.objects.count(),
        'total_kategori': Kategori.objects.count(),
    }

@login_required
def upload_arsip(request):
    if not hasattr(request.user, 'profile') or request.user.profile.role not in ['Admin', 'Pembuat']:
        messages.error(request, 'Anda tidak memiliki izin untuk mengupload arsip.')
        return redirect('arsip:daftar_arsip')
    if request.method == 'POST':
        form = ArsipForm(request.POST, request.FILES)
        if form.is_valid():
            arsip = form.save(commit=False)
            arsip.dibuat_oleh = request.user
            arsip.tanggal_diterima = timezone.now().date()
            arsip.status_verifikasi = 'pending'
            arsip.save()
            form.save_m2m()
            messages.success(request, 'Arsip berhasil diupload dan menunggu verifikasi.')
            return redirect('arsip:detail_arsip', arsip_id=arsip.id)
    else:
        form = ArsipForm()
    return render(request, 'arsip/upload.html', {'form': form})

def register(request):
    if request.method == 'POST':
        form = RegisterForm(request.POST)
        if form.is_valid():
            if form.cleaned_data['password1'] != form.cleaned_data['password2']:
                messages.error(request, "Password tidak cocok!")
                return render(request, 'arsip/register.html', {'form': form})
            user = form.save()
            login(request, user)
            messages.success(request, 'Akun berhasil dibuat.')
            return redirect('arsip:beranda')
    else:
        form = RegisterForm()
    return render(request, 'arsip/register.html', {'form': form})

def password_reset_request(request):
    if request.method == 'POST':
        email = request.POST.get('email', '')
        try:
            user = User.objects.get(email=email)
            return redirect('arsip:password_change', user_id=user.id)
        except User.DoesNotExist:
            messages.error(request, 'Email tidak terdaftar.')
    return render(request, 'arsip/password_reset.html')

def password_change(request, user_id):
    try:
        user = User.objects.get(id=user_id)
    except User.DoesNotExist:
        messages.error(request, 'User tidak ditemukan.')
        return redirect('arsip:password_reset')
    if request.method == 'POST':
        form = PasswordChangeForm(user, request.POST)
        if form.is_valid():
            form.save()
            messages.success(request, 'Password berhasil diubah.')
            return redirect('arsip:login')
    else:
        form = PasswordChangeForm(user)
    return render(request, 'arsip/password_change.html', {'form': form})
